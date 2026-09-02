"""Presentation worker: PPTX native extraction + Vision AI for charts/figures."""

import logging
from io import BytesIO
from pathlib import Path
from typing import Optional

from vision_ai.workers.base import VisionWorker, AnalysisResult, parse_structured_response, AttachmentGroup, PageContext
from vision_ai.client import api_call_with_retry
from vision_ai.vision_utils import pdf_to_images, pdf_page_count, encode_image, build_vision_message

logger = logging.getLogger(__name__)

# Slide-text budget for the prompt. 60,000 chars (~15k tokens) covers
# table-heavy decks with grouped shapes without exhausting the context.
MAX_PROMPT_CHARS = 60000

PPTX_PROMPT = """Analyze this slide content extracted from a presentation.

Slide content:
{slide_text}

Context:
- Page: {page_title}
- Section: {section_path}

Respond in EXACTLY this format (fill in each field):

TITLE: [presentation title]
AUTHORS: [comma-separated list, or "Unknown"]
DATE: [date if mentioned, or "Unknown"]
KEY_POINTS:
- [point 1]
- [point 2]
- [point 3]
BODY:
[2-3 sentence summary of what this presentation covers and its main message]"""

PDF_VISION_PROMPT = """Analyze these presentation slides.

Context:
- Page: {page_title}
- Section: {section_path}

Respond in EXACTLY this format (fill in each field):

TITLE: [presentation title visible on slides]
AUTHORS: [comma-separated list, or "Unknown"]
DATE: [date if visible, or "Unknown"]
KEY_POINTS:
- [key finding/message 1]
- [key finding/message 2]
- [key finding/message 3]
BODY:
[2-3 sentence summary of what these slides cover and the main takeaway]"""


class PresentationWorker(VisionWorker):

    def analyze(self, group: AttachmentGroup, images: dict[str, bytes],
                page_context: PageContext) -> Optional[AnalysisResult]:
        filename = group.filenames[0]
        data = images[filename]
        ext = Path(filename).suffix.lower()

        if ext == '.pptx':
            return self._analyze_pptx(data, filename, page_context)
        else:
            return self._analyze_pdf_presentation(data, filename, page_context)

    def _analyze_pptx(self, data: bytes, filename: str, ctx: PageContext) -> Optional[AnalysisResult]:
        text = self._extract_pptx_text(data)
        if not text:
            return None

        slide_count = text.count('--- Slide ')
        truncated = text[:MAX_PROMPT_CHARS]
        if len(text) > MAX_PROMPT_CHARS:
            truncated += "\n[... truncated]"
            logger.warning(
                f"{filename}: slide text truncated from {len(text)} to "
                f"{MAX_PROMPT_CHARS} chars; summary covers the first "
                f"{truncated.count('--- Slide ')} of {slide_count} slides"
            )

        prompt = PPTX_PROMPT.format(
            slide_text=truncated,
            page_title=ctx.page_title,
            section_path=ctx.section_path,
        )

        messages = [{"role": "user", "content": prompt}]
        try:
            response = api_call_with_retry(messages, max_tokens=2048)
            result = parse_structured_response(response, default_content_type='presentation')
            result.extra['slide_count'] = slide_count
            return result
        except Exception as e:
            logger.error(f"API call failed for presentation {filename}: {e}", exc_info=True)
            return None

    def _shape_texts(self, shape) -> list[str]:
        """Text from one shape, recursing into groups and tables.

        Reading only top-level `has_text_frame` shapes misses three common
        cases: grouped shapes (any diagram assembled from parts), tables (most
        of the text on a comparison or selection slide), and SmartArt/charts
        converted to groups. A table-heavy deck would otherwise extract as
        empty and be skipped with no analysis.
        """
        texts: list[str] = []
        # MSO_SHAPE_TYPE.GROUP == 6; compared numerically to avoid importing
        # the enum, which python-pptx exposes at a version-dependent path.
        if getattr(shape, "shape_type", None) == 6:
            for child in shape.shapes:
                texts.extend(self._shape_texts(child))
            return texts
        if getattr(shape, "has_text_frame", False):
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    texts.append(text)
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                cells = [c.text.strip().replace("\n", " ") for c in row.cells]
                if any(cells):
                    texts.append(" | ".join(cells))
        return texts

    def _extract_pptx_text(self, data: bytes) -> str:
        try:
            from pptx import Presentation
            prs = Presentation(BytesIO(data))
            slides_text = []
            for i, slide in enumerate(prs.slides, 1):
                texts = []
                for shape in slide.shapes:
                    try:
                        texts.extend(self._shape_texts(shape))
                    except Exception as e:
                        # One malformed shape must not cost the whole deck.
                        logger.warning(f"slide {i}: shape skipped ({e})")
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame
                    note_text = notes.text.strip() if notes is not None else ""
                    if note_text:
                        texts.append(f"[speaker notes] {note_text}")
                if texts:
                    slides_text.append(f"--- Slide {i} ---\n" + '\n'.join(texts))
            if not slides_text:
                logger.warning(
                    f"pptx contained no extractable text across "
                    f"{len(prs.slides)} slide(s): likely an image-only deck; "
                    f"no analysis will be written"
                )
            return '\n\n'.join(slides_text)
        except Exception as e:
            logger.warning(f"python-pptx extraction failed: {e}", exc_info=True)
            return ""

    def _analyze_pdf_presentation(self, data: bytes, filename: str,
                                  ctx: PageContext) -> Optional[AnalysisResult]:
        try:
            total = pdf_page_count(data)
        except Exception as e:
            logger.error(f"Failed to read PDF {filename}: {e}")
            return None

        if total == 0:
            return None

        if total <= 8:
            sample_indices = list(range(total))
        else:
            step = total / 8
            sample_indices = [int(i * step) for i in range(8)]

        try:
            pil_images = pdf_to_images(data, dpi=150, page_indices=sample_indices)
        except Exception as e:
            logger.error(f"Failed to render PDF {filename}: {e}")
            return None

        encoded = []
        for img in pil_images:
            encoded.append(encode_image(img, max_dim=1568))

        prompt = f"These are {len(encoded)} sampled slides from a {total}-slide presentation.\n\n"
        prompt += PDF_VISION_PROMPT.format(
            page_title=ctx.page_title,
            section_path=ctx.section_path,
        )

        messages = build_vision_message(encoded, prompt)
        try:
            response = api_call_with_retry(messages, max_tokens=2048)
            result = parse_structured_response(response, default_content_type='presentation')
            result.extra['slide_count'] = total
            return result
        except Exception as e:
            logger.error(f"Vision API failed for {filename}: {e}")
            return None
